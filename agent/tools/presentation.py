from pathlib import Path
from typing import Optional
import re

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# =============================================================
# Configuration
# =============================================================

OUTPUT_DIR = Path(
    "outputs/generated_reports"
)

POWERPOINT_MIME_TYPE = (
    "application/vnd.openxmlformats-officedocument."
    "presentationml.presentation"
)

# Presentation theme
NAVY = RGBColor(13, 35, 64)
BLUE = RGBColor(32, 91, 155)
GOLD = RGBColor(230, 177, 77)

LIGHT_BG = RGBColor(246, 248, 251)
WHITE = RGBColor(255, 255, 255)

DARK_TEXT = RGBColor(30, 41, 59)
MUTED_TEXT = RGBColor(100, 116, 139)

GREEN = RGBColor(22, 163, 74)
RED = RGBColor(220, 38, 38)

BORDER = RGBColor(226, 232, 240)


# =============================================================
# Utility functions
# =============================================================


def _safe_filename(
    value: str,
) -> str:
    """
    Convert a title or filename into a safe filesystem name.
    """

    value = value.strip().lower()

    value = re.sub(
        r"[^a-z0-9]+",
        "_",
        value,
    )

    value = value.strip("_")

    return (
        value
        or "financial_analysis"
    )


def _blank_slide(
    presentation: Presentation,
):
    """
    Add and return a blank slide.
    """

    return presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )


def _add_background(
    slide,
    color: RGBColor,
) -> None:
    """
    Add a solid background rectangle covering the slide.
    """

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )

    background.fill.solid()

    background.fill.fore_color.rgb = (
        color
    )

    background.line.fill.background()


def _add_header(
    slide,
    title: str,
    subtitle: Optional[str] = None,
    *,
    dark_background: bool = False,
) -> None:
    """
    Add a consistent presentation header.
    """

    title_box = slide.shapes.add_textbox(
        Inches(0.65),
        Inches(0.35),
        Inches(11.9),
        Inches(0.65),
    )

    paragraph = (
        title_box
        .text_frame
        .paragraphs[0]
    )

    paragraph.text = title

    paragraph.font.size = Pt(
        26
    )

    paragraph.font.bold = True

    paragraph.font.color.rgb = (
        WHITE
        if dark_background
        else NAVY
    )

    if subtitle:
        subtitle_box = (
            slide.shapes.add_textbox(
                Inches(0.68),
                Inches(0.98),
                Inches(11.4),
                Inches(0.35),
            )
        )

        paragraph = (
            subtitle_box
            .text_frame
            .paragraphs[0]
        )

        paragraph.text = subtitle

        paragraph.font.size = Pt(
            12
        )

        paragraph.font.color.rgb = (
            RGBColor(
                210,
                220,
                232,
            )
            if dark_background
            else MUTED_TEXT
        )


def _add_footer(
    slide,
    source: Optional[str] = None,
    *,
    dark_background: bool = False,
) -> None:
    """
    Add a bottom accent line and optional source citation.
    """

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(7.12),
        Inches(13.333),
        Inches(0.04),
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    if not source:
        return

    source_box = slide.shapes.add_textbox(
        Inches(0.65),
        Inches(7.19),
        Inches(12.0),
        Inches(0.22),
    )

    paragraph = (
        source_box
        .text_frame
        .paragraphs[0]
    )

    paragraph.text = (
        f"Source: {source}"
    )

    paragraph.font.size = Pt(
        8.5
    )

    paragraph.font.color.rgb = (
        RGBColor(
            205,
            215,
            225,
        )
        if dark_background
        else MUTED_TEXT
    )


# =============================================================
# Title slide
# =============================================================


def _add_title_slide(
    presentation: Presentation,
    title: str,
    subtitle: Optional[str] = None,
) -> None:
    """
    Add the main title slide.
    """

    slide = _blank_slide(
        presentation
    )

    _add_background(
        slide,
        NAVY,
    )

    # Decorative accent
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(6.85),
        Inches(13.333),
        Inches(0.18),
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.85),
        Inches(2.2),
        Inches(11.5),
        Inches(1.35),
    )

    paragraph = (
        title_box
        .text_frame
        .paragraphs[0]
    )

    paragraph.text = title

    paragraph.font.size = Pt(
        40
    )

    paragraph.font.bold = True
    paragraph.font.color.rgb = WHITE

    if subtitle:
        subtitle_box = (
            slide.shapes.add_textbox(
                Inches(0.9),
                Inches(3.65),
                Inches(10.8),
                Inches(0.75),
            )
        )

        paragraph = (
            subtitle_box
            .text_frame
            .paragraphs[0]
        )

        paragraph.text = subtitle

        paragraph.font.size = Pt(
            18
        )

        paragraph.font.color.rgb = (
            RGBColor(
                220,
                230,
                242,
            )
        )


# =============================================================
# Standard bullet slide
# =============================================================


def _add_bullet_slide(
    presentation: Presentation,
    title: str,
    bullets: list[str],
    source: Optional[str] = None,
) -> None:
    """
    Add a clean executive-style bullet slide.
    """

    slide = _blank_slide(
        presentation
    )

    _add_background(
        slide,
        WHITE,
    )

    _add_header(
        slide,
        title,
    )

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        Inches(1.3),
        Inches(0.08),
        Inches(4.9),
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    text_box = slide.shapes.add_textbox(
        Inches(1.05),
        Inches(1.45),
        Inches(11.1),
        Inches(5.1),
    )

    text_frame = (
        text_box.text_frame
    )

    text_frame.word_wrap = True

    safe_bullets = (
        bullets
        if bullets
        else [
            "No additional information provided."
        ]
    )

    for index, bullet in enumerate(
        safe_bullets
    ):
        if index == 0:
            paragraph = (
                text_frame.paragraphs[0]
            )

        else:
            paragraph = (
                text_frame.add_paragraph()
            )

        paragraph.text = str(
            bullet
        )

        paragraph.font.size = Pt(
            19
        )

        paragraph.font.color.rgb = (
            DARK_TEXT
        )

        paragraph.space_after = Pt(
            12
        )

        paragraph.level = 0

    _add_footer(
        slide,
        source,
    )


# =============================================================
# Metric cards
# =============================================================


def _add_metric_card(
    slide,
    *,
    x: float,
    y: float,
    width: float,
    height: float,
    label: str,
    value: str,
    note: Optional[str] = None,
    accent_color: RGBColor = BLUE,
) -> None:
    """
    Add a financial KPI card.
    """

    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.05),
        Inches(y + 0.06),
        Inches(width),
        Inches(height),
    )

    shadow.fill.solid()

    shadow.fill.fore_color.rgb = (
        RGBColor(
            230,
            233,
            238,
        )
    )

    shadow.line.fill.background()

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )

    card.fill.solid()
    card.fill.fore_color.rgb = WHITE

    card.line.color.rgb = BORDER

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.08),
        Inches(height),
    )

    accent.fill.solid()

    accent.fill.fore_color.rgb = (
        accent_color
    )

    accent.line.fill.background()

    label_box = slide.shapes.add_textbox(
        Inches(x + 0.3),
        Inches(y + 0.28),
        Inches(width - 0.55),
        Inches(0.35),
    )

    paragraph = (
        label_box
        .text_frame
        .paragraphs[0]
    )

    paragraph.text = label.upper()

    paragraph.font.size = Pt(
        11
    )

    paragraph.font.bold = True

    paragraph.font.color.rgb = (
        MUTED_TEXT
    )

    value_box = slide.shapes.add_textbox(
        Inches(x + 0.3),
        Inches(y + 0.8),
        Inches(width - 0.55),
        Inches(0.7),
    )

    paragraph = (
        value_box
        .text_frame
        .paragraphs[0]
    )

    paragraph.text = value

    paragraph.font.size = Pt(
        25
    )

    paragraph.font.bold = True

    paragraph.font.color.rgb = (
        accent_color
    )

    if note:
        note_box = (
            slide.shapes.add_textbox(
                Inches(x + 0.3),
                Inches(y + 1.55),
                Inches(width - 0.55),
                Inches(0.4),
            )
        )

        paragraph = (
            note_box
            .text_frame
            .paragraphs[0]
        )

        paragraph.text = note

        paragraph.font.size = Pt(
            10
        )

        paragraph.font.color.rgb = (
            MUTED_TEXT
        )


def _add_metrics_slide(
    presentation: Presentation,
    title: str,
    metrics: list[dict],
    source: Optional[str] = None,
) -> None:
    """
    Add a KPI / financial metrics slide.
    """

    slide = _blank_slide(
        presentation
    )

    _add_background(
        slide,
        LIGHT_BG,
    )

    _add_header(
        slide,
        title,
    )

    metrics = metrics[:4]

    if not metrics:
        metrics = [
            {
                "label": "Metric",
                "value": "N/A",
            }
        ]

    count = len(metrics)

    if count == 1:
        width = 5.5
        start_x = 3.9

    elif count == 2:
        width = 5.3
        start_x = 1.05

    elif count == 3:
        width = 3.75
        start_x = 0.65

    else:
        width = 2.85
        start_x = 0.55

    gap = 0.35

    for index, metric in enumerate(
        metrics
    ):
        x = (
            start_x
            + index
            * (
                width
                + gap
            )
        )

        accent_color = (
            GOLD
            if index == count - 1
            else BLUE
        )

        _add_metric_card(
            slide,
            x=x,
            y=2.15,
            width=width,
            height=2.35,
            label=str(
                metric.get(
                    "label",
                    "",
                )
            ),
            value=str(
                metric.get(
                    "value",
                    "",
                )
            ),
            note=metric.get(
                "note"
            ),
            accent_color=(
                accent_color
            ),
        )

    _add_footer(
        slide,
        source,
    )


# =============================================================
# Highlight slide
# =============================================================


def _add_highlight_slide(
    presentation: Presentation,
    title: str,
    headline: str,
    subtitle: Optional[str] = None,
    source: Optional[str] = None,
) -> None:
    """
    Add a large-number highlight slide.
    """

    slide = _blank_slide(
        presentation
    )

    _add_background(
        slide,
        NAVY,
    )

    _add_header(
        slide,
        title,
        dark_background=True,
    )

    headline_box = (
        slide.shapes.add_textbox(
            Inches(0.8),
            Inches(2.2),
            Inches(11.75),
            Inches(1.5),
        )
    )

    paragraph = (
        headline_box
        .text_frame
        .paragraphs[0]
    )

    paragraph.text = str(
        headline
    )

    paragraph.font.size = Pt(
        58
    )

    paragraph.font.bold = True
    paragraph.font.color.rgb = GOLD

    paragraph.alignment = (
        PP_ALIGN.CENTER
    )

    if subtitle:
        subtitle_box = (
            slide.shapes.add_textbox(
                Inches(1.4),
                Inches(3.85),
                Inches(10.5),
                Inches(0.9),
            )
        )

        paragraph = (
            subtitle_box
            .text_frame
            .paragraphs[0]
        )

        paragraph.text = str(
            subtitle
        )

        paragraph.font.size = Pt(
            21
        )

        paragraph.font.color.rgb = WHITE

        paragraph.alignment = (
            PP_ALIGN.CENTER
        )

    _add_footer(
        slide,
        source,
        dark_background=True,
    )


# =============================================================
# Comparison slide
# =============================================================


def _add_comparison_slide(
    presentation: Presentation,
    title: str,
    left: dict,
    right: dict,
    conclusion: Optional[str] = None,
    source: Optional[str] = None,
) -> None:
    """
    Add a before-vs-after / year-vs-year comparison slide.
    """

    slide = _blank_slide(
        presentation
    )

    _add_background(
        slide,
        LIGHT_BG,
    )

    _add_header(
        slide,
        title,
    )

    _add_metric_card(
        slide,
        x=0.95,
        y=1.95,
        width=5.35,
        height=2.35,
        label=str(
            left.get(
                "label",
                "Previous",
            )
        ),
        value=str(
            left.get(
                "value",
                "",
            )
        ),
        note=left.get(
            "note"
        ),
        accent_color=BLUE,
    )

    _add_metric_card(
        slide,
        x=7.0,
        y=1.95,
        width=5.35,
        height=2.35,
        label=str(
            right.get(
                "label",
                "Current",
            )
        ),
        value=str(
            right.get(
                "value",
                "",
            )
        ),
        note=right.get(
            "note"
        ),
        accent_color=GOLD,
    )

    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(6.2),
        Inches(2.7),
        Inches(0.62),
        Inches(0.42),
    )

    arrow.fill.solid()

    arrow.fill.fore_color.rgb = (
        MUTED_TEXT
    )

    arrow.line.fill.background()

    if conclusion:
        conclusion_box = (
            slide.shapes.add_textbox(
                Inches(1.0),
                Inches(5.05),
                Inches(11.3),
                Inches(0.95),
            )
        )

        paragraph = (
            conclusion_box
            .text_frame
            .paragraphs[0]
        )

        paragraph.text = str(
            conclusion
        )

        paragraph.font.size = Pt(
            21
        )

        paragraph.font.bold = True
        paragraph.font.color.rgb = NAVY

        paragraph.alignment = (
            PP_ALIGN.CENTER
        )

    _add_footer(
        slide,
        source,
    )


# =============================================================
# Chart slide
# =============================================================


def _add_chart_slide(
    presentation: Presentation,
    title: str,
    categories: list,
    values: list,
    source: Optional[str] = None,
) -> None:
    """
    Add a simple column-chart slide.
    """

    if not categories:
        raise ValueError(
            "Chart requires at least one category."
        )

    if len(categories) != len(values):
        raise ValueError(
            "Chart categories and values must have equal lengths."
        )

    slide = _blank_slide(
        presentation
    )

    _add_background(
        slide,
        WHITE,
    )

    _add_header(
        slide,
        title,
    )

    chart_data = (
        CategoryChartData()
    )

    chart_data.categories = [
        str(category)
        for category in categories
    ]

    numeric_values = [
        float(value)
        for value in values
    ]

    chart_data.add_series(
        "Value",
        numeric_values,
    )

    chart_shape = (
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1.0),
            Inches(1.55),
            Inches(11.4),
            Inches(4.95),
            chart_data,
        )
    )

    chart = chart_shape.chart

    chart.has_legend = False

    chart.value_axis.has_major_gridlines = (
        True
    )

    chart.category_axis.tick_labels.font.size = (
        Pt(11)
    )

    chart.value_axis.tick_labels.font.size = (
        Pt(10)
    )

    series = chart.series[0]

    series.format.fill.solid()

    series.format.fill.fore_color.rgb = (
        BLUE
    )

    _add_footer(
        slide,
        source,
    )


# =============================================================
# Sources slide
# =============================================================


def _add_sources_slide(
    presentation: Presentation,
    title: str,
    sources: list[str],
) -> None:
    """
    Add a dedicated references slide.
    """

    slide = _blank_slide(
        presentation
    )

    _add_background(
        slide,
        WHITE,
    )

    _add_header(
        slide,
        title,
    )

    text_box = slide.shapes.add_textbox(
        Inches(0.85),
        Inches(1.45),
        Inches(11.6),
        Inches(5.5),
    )

    text_frame = (
        text_box.text_frame
    )

    text_frame.word_wrap = True

    safe_sources = (
        sources
        if sources
        else [
            "No sources supplied."
        ]
    )

    for index, source in enumerate(
        safe_sources
    ):
        if index == 0:
            paragraph = (
                text_frame.paragraphs[0]
            )

        else:
            paragraph = (
                text_frame.add_paragraph()
            )

        paragraph.text = str(
            source
        )

        paragraph.font.size = Pt(
            14
        )

        paragraph.font.color.rgb = (
            DARK_TEXT
        )

        paragraph.space_after = Pt(
            9
        )

    _add_footer(
        slide
    )


# =============================================================
# Public Gemini tool
# =============================================================


def create_powerpoint(
    title: str,
    slides: list[dict],
    subtitle: Optional[str] = None,
    filename: Optional[str] = None,
) -> dict:
    """
    Create a designed PowerPoint presentation and save it in
    the generated reports directory.

    Use this tool only when the user explicitly requests a
    PowerPoint, presentation, slide deck, or slides.

    Supported slide types are:

    bullets:
        {
            "type": "bullets",
            "title": "Executive Summary",
            "bullets": [
                "First insight",
                "Second insight"
            ],
            "source": "GTCO 2024 Annual Report, p. 247"
        }

    metrics:
        {
            "type": "metrics",
            "title": "Key Financial Metrics",
            "metrics": [
                {
                    "label": "2023 PBT",
                    "value": "₦608.8bn",
                    "note": "Group"
                },
                {
                    "label": "2024 PBT",
                    "value": "₦1.266tn",
                    "note": "Group"
                }
            ],
            "source": "GTCO 2024 Annual Report"
        }

    highlight:
        {
            "type": "highlight",
            "title": "Profit Growth",
            "headline": "107.98%",
            "subtitle": "Increase in profit before tax",
            "source": "Calculated from annual-report values"
        }

    comparison:
        {
            "type": "comparison",
            "title": "Profit Before Tax",
            "left": {
                "label": "2023",
                "value": "₦608.8bn"
            },
            "right": {
                "label": "2024",
                "value": "₦1.266tn"
            },
            "conclusion": "Profit before tax more than doubled.",
            "source": "GTCO Annual Report"
        }

    chart:
        {
            "type": "chart",
            "title": "Profit Before Tax Trend",
            "categories": [
                "2023",
                "2024"
            ],
            "values": [
                608.8,
                1266.2
            ],
            "source": "Values shown in ₦bn"
        }

    sources:
        {
            "type": "sources",
            "title": "Sources",
            "sources": [
                "GTCO 2024 Annual Report, p. 247"
            ]
        }

    Args:
        title:
            Main presentation title.

        slides:
            Ordered list of slide configuration objects.

        subtitle:
            Optional presentation subtitle.

        filename:
            Optional output filename without a directory path.

    Returns:
        Information about the generated PowerPoint file,
        including the filename and metadata required by the API
        to make the presentation downloadable.
    """

    try:
        if not slides:
            return {
                "success": False,
                "error_type": (
                    "ValueError"
                ),
                "error": (
                    "At least one content "
                    "slide is required."
                ),
            }

        OUTPUT_DIR.mkdir(
            parents=True,
            exist_ok=True,
        )

        raw_name = (
            Path(filename).stem
            if filename
            else title
        )

        safe_name = _safe_filename(
            raw_name
        )

        output_path = (
            OUTPUT_DIR
            / f"{safe_name}.pptx"
        )

        presentation = Presentation()

        # 16:9 widescreen
        presentation.slide_width = (
            Inches(13.333)
        )

        presentation.slide_height = (
            Inches(7.5)
        )

        # =====================================================
        # Title slide
        # =====================================================

        _add_title_slide(
            presentation,
            title,
            subtitle,
        )

        # =====================================================
        # Content slides
        # =====================================================

        for slide_data in slides:
            if not isinstance(
                slide_data,
                dict,
            ):
                continue

            slide_type = str(
                slide_data.get(
                    "type",
                    "bullets",
                )
            ).lower()

            slide_title = str(
                slide_data.get(
                    "title",
                    "Untitled",
                )
            )

            source = slide_data.get(
                "source"
            )

            if slide_type == "metrics":
                _add_metrics_slide(
                    presentation,
                    slide_title,
                    slide_data.get(
                        "metrics",
                        [],
                    ),
                    source,
                )

            elif slide_type == "highlight":
                _add_highlight_slide(
                    presentation,
                    slide_title,
                    str(
                        slide_data.get(
                            "headline",
                            "",
                        )
                    ),
                    slide_data.get(
                        "subtitle"
                    ),
                    source,
                )

            elif slide_type == "comparison":
                _add_comparison_slide(
                    presentation,
                    slide_title,
                    slide_data.get(
                        "left",
                        {},
                    ),
                    slide_data.get(
                        "right",
                        {},
                    ),
                    slide_data.get(
                        "conclusion"
                    ),
                    source,
                )

            elif slide_type == "chart":
                _add_chart_slide(
                    presentation,
                    slide_title,
                    slide_data.get(
                        "categories",
                        [],
                    ),
                    slide_data.get(
                        "values",
                        [],
                    ),
                    source,
                )

            elif slide_type == "sources":
                _add_sources_slide(
                    presentation,
                    slide_title,
                    slide_data.get(
                        "sources",
                        [],
                    ),
                )

            else:
                _add_bullet_slide(
                    presentation,
                    slide_title,
                    slide_data.get(
                        "bullets",
                        [],
                    ),
                    source,
                )

        # =====================================================
        # Save
        # =====================================================

        presentation.save(
            output_path
        )

        return {
            "success": True,
            "title": title,
            "format": "pptx",
            "file_type": "powerpoint",
            "filename": output_path.name,
            "path": str(output_path),
            "mime_type": POWERPOINT_MIME_TYPE,
            "size_bytes": (
                output_path.stat().st_size
            ),
            "slides_created": len(
                presentation.slides
            ),
        }

    except Exception as exc:
        return {
            "success": False,
            "error_type": (
                type(exc).__name__
            ),
            "error": str(exc),
        }